"""会議用資料の出力（proposal.md 2-3）。

pdf / word / 画像入り pptx を作る。いずれも
  - 章 = 案件、節 = 段階、項 = 記録
という構成に揃え、各節の冒頭に slides.py が作った資料用画像を貼る。

既存の exporters.py の関数は書き換えない。ここは追加の builder で、
exporters.BUILDERS へ後から登録する。
"""

from __future__ import annotations

import io
import sys
from datetime import datetime
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from ledger import Ledger  # noqa: E402
from projecttree import slides as _slides  # noqa: E402

STAGE_TITLES = _slides.STAGE_TITLES


def _data(ledger: Ledger, thread_id: str, stage_no: int | None):
    """exporters の収集関数を再利用する（重複実装を作らない）。"""
    from projecttree import exporters as _exp
    return _exp._collect(ledger, thread_id, stage_no)


def _stage_nums(stage_no: int | None) -> list[int]:
    return [stage_no] if stage_no else [1, 2, 3, 4, 5]


def _stage_images(ledger: Ledger, thread_id: str) -> dict[int, Path]:
    """段階ごとに資料へ貼る画像を返す（Enhancement02.md 1-1）。

    イメージ図（illustrate.py が LLM に描かせた二次画像）があればそれを使う。
    無い段階は、これまで通り slides.py の資料用画像で埋める。
    どちらも「段階番号 → 画像パス」という同じ形なので、貼る側は区別しなくてよい。
    """
    imgs = _slides.ensure_images(ledger, thread_id)
    try:
        from projecttree import illustrate as _ill
        for n, p in _ill.image_paths(ledger, thread_id).items():
            imgs[n] = p
    except Exception:
        pass          # イメージ図が無くても資料は出せる。ここで資料出力を止めない。
    return imgs


# --------------------------------------------------------------------------
# PDF（PyMuPDF）
# --------------------------------------------------------------------------

# 組込みの "japan" は ASCII まで全角幅で組むため、数字と日付が読みにくい。
# 実フォントを埋め込み、保存時にサブセット化する。
_JP_FONT = "japan"
_JP_ALIAS = "jpfont"


def _jp_fontfile() -> str | None:
    for p in _slides._FONT_CANDIDATES:
        if Path(p).exists() and p.lower().endswith((".ttf", ".ttc")):
            return p
    return None


def _png_as_jpeg(path: Path, quality: int = 78) -> bytes:
    """PDF へ貼る用に JPEG 化する。PNG のまま貼ると PDF が数MB に膨らむ。"""
    from PIL import Image
    buf = io.BytesIO()
    Image.open(path).convert("RGB").save(buf, "JPEG", quality=quality, optimize=True)
    return buf.getvalue()


def build_pdf(ledger: Ledger, thread_id: str, stage_no: int | None = None) -> bytes:
    import fitz

    d = _data(ledger, thread_id, stage_no)
    imgs = _stage_images(ledger, thread_id)
    doc = fitz.open()
    PW, PH = 595, 842            # A4 縦（pt）
    M = 50

    fontfile = _jp_fontfile()
    fname = _JP_ALIAS if fontfile else _JP_FONT
    # 幅の計測は fitz.get_text_length では埋め込みフォントを扱えないので Font 側で行う。
    measure = fitz.Font(fontfile=fontfile) if fontfile else None

    def width(s: str, size: float) -> float:
        if measure:
            return measure.text_length(s, fontsize=size)
        return fitz.get_text_length(s, fontname=_JP_FONT, fontsize=size)

    def new_page():
        p = doc.new_page(width=PW, height=PH)
        if fontfile:
            p.insert_font(fontname=_JP_ALIAS, fontfile=fontfile)
        return p, M

    def text(page, y, s, size=10.5, color=(0.1, 0.1, 0.1), indent=0):
        """折り返しながら書き、次の y を返す。"""
        maxw = PW - M * 2 - indent
        cur, out = "", []
        for ch in s:
            w = width(cur + ch, size)
            if w > maxw and cur:
                out.append(cur); cur = ch
            else:
                cur += ch
        if cur:
            out.append(cur)
        for ln in out:
            page.insert_text((M + indent, y), ln, fontname=fname, fontsize=size, color=color)
            y += size * 1.55
        return y

    # 表紙相当の見出し（章）
    page, y = new_page()
    y = text(page, y + 10, d["thread"]["name"], size=19)
    y = text(page, y + 4, f"会議用資料 / 出力日 {datetime.now().strftime('%Y-%m-%d')}",
             size=9.5, color=(0.42, 0.45, 0.5))
    page.draw_line(fitz.Point(M, y), fitz.Point(PW - M, y), color=(0.75, 0.77, 0.8), width=0.7)
    y += 18

    for n in _stage_nums(stage_no):
        st = next((s for s in d["stages"] if s["stage_no"] == n), None)
        evs = [e for e in d["events"] if e.get("stage_no") == n]

        # 節ごとにページを改める（画像＋本文で1ページに収まる想定）
        if y > PH - 330:
            page, y = new_page()

        rgb = tuple(c / 255 for c in _slides.STAGE_COLOR[n])
        page.draw_rect(fitz.Rect(M, y + 2, M + 4, y + 20), color=rgb, fill=rgb)
        y = text(page, y + 15, f"{n}. {STAGE_TITLES[n]}", size=14)
        y += 4

        # 資料用画像（スライドと同じ絵をそのまま図版として使う）
        if n in imgs:
            iw = PW - M * 2
            ih = iw * _slides.H / _slides.W
            if y + ih > PH - M:
                page, y = new_page()
            page.insert_image(fitz.Rect(M, y, M + iw, y + ih), stream=_png_as_jpeg(imgs[n]))
            y += ih + 12

        if st and st.get("summary"):
            y = text(page, y, st["summary"], size=10.5)
            y += 6

        for e in evs:
            if y > PH - M - 24:
                page, y = new_page()
            y = text(page, y, f"・{e['occurred_on']}  {e['summary']}", size=10, indent=6)
        y += 12

    if d.get("gaps"):
        if y > PH - 160:
            page, y = new_page()
        y = text(page, y + 6, "記録の空白（SQL が検出）", size=12)
        for g in d["gaps"]:
            if y > PH - M - 24:
                page, y = new_page()
            y = text(page, y, f"・{g.get('description') or g.get('kind')}", size=10, indent=6)

    if fontfile:
        doc.subset_fonts()          # 使った字だけ残す。埋め込みでも数百KBに収まる。
    out = doc.tobytes(garbage=4, deflate=True)
    doc.close()
    return out


# --------------------------------------------------------------------------
# Word（python-docx）
# --------------------------------------------------------------------------

def build_docx(ledger: Ledger, thread_id: str, stage_no: int | None = None) -> bytes:
    from docx import Document
    from docx.shared import Inches, Pt

    d = _data(ledger, thread_id, stage_no)
    imgs = _stage_images(ledger, thread_id)

    doc = Document()
    doc.add_heading(d["thread"]["name"], level=0)
    p = doc.add_paragraph(f"会議用資料 / 出力日 {datetime.now().strftime('%Y-%m-%d')}")
    p.runs[0].font.size = Pt(9)

    for n in _stage_nums(stage_no):
        st = next((s for s in d["stages"] if s["stage_no"] == n), None)
        evs = [e for e in d["events"] if e.get("stage_no") == n]

        doc.add_heading(f"{n}. {STAGE_TITLES[n]}", level=1)
        if n in imgs:
            doc.add_picture(str(imgs[n]), width=Inches(6.3))
        if st and st.get("summary"):
            doc.add_paragraph(st["summary"])
        if evs:
            doc.add_heading("記録", level=2)
            for e in evs:
                doc.add_paragraph(f"{e['occurred_on']}  {e['summary']}", style="List Bullet")

    if d.get("gaps"):
        doc.add_heading("記録の空白（SQL が検出）", level=1)
        for g in d["gaps"]:
            doc.add_paragraph(g.get("description") or g.get("kind"), style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------
# 画像入り pptx（既存の build_pptx は残し、別形式として追加）
# --------------------------------------------------------------------------

def build_pptx_slides(ledger: Ledger, thread_id: str, stage_no: int | None = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    d = _data(ledger, thread_id, stage_no)
    imgs = _stage_images(ledger, thread_id)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title = prs.slides.add_slide(prs.slide_layouts[5])
    title.shapes.title.text = d["thread"]["name"]

    for n in _stage_nums(stage_no):
        s = prs.slides.add_slide(blank)
        if n in imgs:
            # 16:9 の画像をスライド全面に敷く
            s.shapes.add_picture(str(imgs[n]), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        else:
            tb = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(1))
            tb.text_frame.text = f"{n}. {STAGE_TITLES[n]}"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
