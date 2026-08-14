"""1次情報の入力（Enhancement.md §2-1）。

プロンプト・資料・画像を受け取り、台帳の documents として取り込む。
既存の ingest.py と同じ正規化・doc_id 規則を使うため、取り込んだ資料は
そのまま extractor.py の対象になる（1文書につき生涯1回だけ抽出）。

LLM は呼ばない。API 復旧後は extractor を回すだけで段階分類の精度が上がる。
"""

from __future__ import annotations

import hashlib
import io
import json
import re
from datetime import date, datetime, timezone

from ledger import Ledger
from projecttree import security

# Enhancement.md §2-1 の対象資料・対象画像
TEXT_EXT = {".txt", ".md"}
DOC_EXT = {".pdf", ".docx", ".pptx"}
IMAGE_EXT = {".jpg", ".jpeg", ".png"}
ACCEPTED_EXT = TEXT_EXT | DOC_EXT | IMAGE_EXT

MAX_UPLOAD_BYTES = 30 * 1024 * 1024
MAX_TEXT_CHARS = 400_000


class IntakeRejected(Exception):
    """取り込みを拒否した理由を利用者に返すための例外。"""


# ---------------------------------------------------------------- 抽出

def _from_pdf(data: bytes, filename: str) -> str:
    # 既存のセキュリティ検証（マジックナンバー・暗号化/JS/添付の拒否）を通す
    security.validate_pdf(data, filename)
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        return "\n".join(page.get_text() for page in doc)
    finally:
        doc.close()


def _from_docx(data: bytes) -> str:
    # python-docx は未導入のため、docx(=zip) の XML から本文を取り出す
    import xml.etree.ElementTree as ET
    import zipfile

    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        xml = z.read("word/document.xml")
    root = ET.fromstring(xml)
    lines = []
    for p in root.iter(f"{{{ns['w']}}}p"):
        text = "".join(t.text or "" for t in p.iter(f"{{{ns['w']}}}t"))
        if text.strip():
            lines.append(text)
    return "\n".join(lines)


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## スライド {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        lines.append(t)
    return "\n".join(lines)


def extract_text(data: bytes, filename: str) -> tuple[str, str]:
    """(本文, source_type) を返す。画像は本文を持たない。"""
    ext = ("." + filename.rsplit(".", 1)[-1]).lower() if "." in filename else ""
    if ext not in ACCEPTED_EXT:
        raise IntakeRejected(f"未対応の形式です: {ext or filename}")
    if len(data) > MAX_UPLOAD_BYTES:
        raise IntakeRejected(f"ファイルが大きすぎます（上限 {MAX_UPLOAD_BYTES // 1024 // 1024}MB）")

    if ext in TEXT_EXT:
        for enc in ("utf-8", "cp932", "utf-16"):
            try:
                return data.decode(enc), "memo"
            except UnicodeDecodeError:
                continue
        raise IntakeRejected("文字コードを判定できませんでした")
    if ext == ".pdf":
        return _from_pdf(data, filename), "report"
    if ext == ".docx":
        return _from_docx(data), "report"
    if ext == ".pptx":
        return _from_pptx(data), "report"
    if ext in IMAGE_EXT:
        # 画像は本文を持たない。台帳には添付として記録し、参照だけ残す。
        return "", "attachment"
    raise IntakeRejected(f"未対応の形式です: {ext}")


# ---------------------------------------------------------------- 取り込み

def _normalize(raw: str) -> str:
    """ingest.py と同一の正規化。ここを変えると doc_id が変わるので触らないこと。"""
    text = raw.replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.rstrip() for line in text.split("\n")]
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    return "\n".join(lines)


def _guess_date(text: str, filename: str) -> str:
    """本文・ファイル名から日付を推定する。取れなければ今日。"""
    m = re.search(r"(20\d{2})[-/年](\d{1,2})[-/月](\d{1,2})", text[:2000])
    if m:
        y, mo, d = (int(g) for g in m.groups())
        try:
            return date(y, mo, d).isoformat()
        except ValueError:
            pass
    m = re.match(r"^(\d{4})(\d{2})(\d{2})_", filename)
    if m:
        try:
            return date(*(int(g) for g in m.groups())).isoformat()
        except ValueError:
            pass
    return date.today().isoformat()


def _title(text: str, filename: str) -> str:
    for line in text.split("\n"):
        s = line.strip().lstrip("#").strip()
        if s:
            return s[:120]
    return filename


def intake_file(ledger: Ledger, data: bytes, filename: str, thread_id: str | None = None) -> dict:
    """1ファイルを台帳に取り込む。既存と同一内容なら追加せずスキップする。"""
    safe = security.safe_name(filename)
    text, source_type = extract_text(data, safe)

    if source_type == "attachment":
        # 画像は本文が無いので、参照情報だけを本文として持たせる
        text = f"[画像ファイル] {safe}\n取込日: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    if len(text) > MAX_TEXT_CHARS:
        raise IntakeRejected(f"本文が長すぎます（{len(text):,} 文字 / 上限 {MAX_TEXT_CHARS:,}）")

    norm = _normalize(text)
    if not norm.strip():
        raise IntakeRejected("本文を抽出できませんでした")

    sha = hashlib.sha256(norm.encode("utf-8")).hexdigest()
    if ledger.document_exists(sha):
        return {"status": "skipped", "reason": "同一内容の資料が既に台帳にあります", "filename": safe}

    doc = {
        "doc_id": "doc_" + sha[:12],
        "content_sha256": sha,
        "source_path": f"intake/{safe}",
        "source_type": source_type,
        "title": _title(norm, safe),
        "occurred_at": _guess_date(norm, safe),
        "ingested_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "author": None,
        "participants": [],
        "text": norm,
        "line_count": len(norm.split("\n")),
    }
    ledger.insert_document(doc)
    return {
        "status": "added",
        "doc_id": doc["doc_id"],
        "filename": safe,
        "title": doc["title"],
        "occurred_at": doc["occurred_at"],
        "source_type": source_type,
        "line_count": doc["line_count"],
    }


def intake_prompt(ledger: Ledger, prompt: str, thread_id: str | None = None) -> dict:
    """プロンプト（メモ書き）を1次情報として台帳に取り込む。"""
    if not prompt.strip():
        raise IntakeRejected("入力が空です")
    stamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    body = f"# 入力メモ（{stamp}）\n\n{prompt.strip()}"
    return intake_file(ledger, body.encode("utf-8"), f"prompt_{datetime.now():%Y%m%d_%H%M%S}.md", thread_id)


def pending_count(ledger: Ledger) -> int:
    """取り込み済みだが未抽出（extractor 未実行）の資料数。"""
    return len(ledger.documents_without_events())
