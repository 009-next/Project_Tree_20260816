"""資料出力（md / pptx / xlsx）。

Enhancement.md §2-3-1 の表に対応する。
  md    … "業務確定台帳"の時系列内容と文章のままで表示した資料
  pptx  … プレゼン資料の一部資料
  xlsx  … その内容に関連するグラフや表の資料

台帳（ledger.db）を読むだけで、LLM は呼ばない。既存モジュールは書き換えず import して使う。
"""

from __future__ import annotations

import io
import json
from datetime import datetime

from ledger import Ledger
from projecttree.stages import STAGE_TITLES

# 段階の意味色（projecttree.html の CSS と対応させる）
_STAGE_RGB = {
    1: (0x37, 0x4A, 0x67),
    2: (0x8D, 0x6E, 0x1F),
    3: (0x1F, 0x5E, 0x8D),
    4: (0x8D, 0x3B, 0x1F),
    5: (0xB0, 0x2A, 0x2A),  # 解決策・発展策 = 赤（Enhancement.md §2-2）
}


# ---------------------------------------------------------------- 台帳の読み出し

def _thread(ledger: Ledger, thread_id: str) -> dict:
    r = ledger.conn.execute(
        "SELECT thread_id, name, first_seen, last_seen FROM threads WHERE thread_id = ?",
        (thread_id,),
    ).fetchone()
    if r is None:
        raise ValueError(f"thread not found: {thread_id}")
    return dict(r)


def _stages(ledger: Ledger, thread_id: str) -> list[dict]:
    rows = ledger.conn.execute(
        "SELECT stage_id, stage_no, title, summary, method FROM stages "
        "WHERE thread_id = ? ORDER BY stage_no",
        (thread_id,),
    ).fetchall()
    return [dict(r) for r in rows]


def _stage_events(ledger: Ledger, thread_id: str, stage_no: int | None = None) -> list[dict]:
    """段階に紐付く時系列。stage_no=None なら全段階分をまとめて返す。"""
    sql = """
        SELECT s.stage_no, e.event_id, e.occurred_on, e.kind, e.summary,
               e.span_quote, e.certainty, d.title AS doc_title, d.source_path
        FROM stage_events se
        JOIN stages s   ON se.stage_id = s.stage_id
        JOIN events e   ON se.event_id = e.event_id
        JOIN documents d ON e.doc_id = d.doc_id
        WHERE s.thread_id = ?
    """
    params: list = [thread_id]
    if stage_no is not None:
        sql += " AND s.stage_no = ?"
        params.append(stage_no)
    sql += " ORDER BY s.stage_no, e.occurred_on"
    return [dict(r) for r in ledger.conn.execute(sql, params).fetchall()]


def _gaps(ledger: Ledger, thread_id: str) -> list[dict]:
    rows = ledger.conn.execute(
        "SELECT kind, period_start, description FROM gaps WHERE thread_id = ? ORDER BY period_start",
        (thread_id,),
    ).fetchall()
    return [dict(r) for r in rows]


def _collect(ledger: Ledger, thread_id: str, stage_no: int | None = None) -> dict:
    return {
        "thread": _thread(ledger, thread_id),
        "stages": _stages(ledger, thread_id),
        "events": _stage_events(ledger, thread_id, stage_no),
        "gaps": _gaps(ledger, thread_id),
        "stage_no": stage_no,
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
    }


def filename(thread_name: str, fmt: str, stage_no: int | None = None) -> str:
    stamp = datetime.now().strftime("%Y%m%d")
    safe = "".join(c for c in thread_name if c not in '\\/:*?"<>|').strip() or "project"
    part = f"_段階{stage_no}" if stage_no else ""
    return f"{safe}{part}_{stamp}.{fmt}"


# ---------------------------------------------------------------- md

def build_md(ledger: Ledger, thread_id: str, stage_no: int | None = None) -> bytes:
    """時系列内容と原文をそのまま章節項に構成した資料（Enhancement.md §2-2）。"""
    d = _collect(ledger, thread_id, stage_no)
    th = d["thread"]
    out: list[str] = [
        f"# {th['name']}",
        "",
        f"- 期間: {th['first_seen']} 〜 {th['last_seen']}",
        f"- 出力日: {d['generated_at']}",
        f"- 記録件数: {len(d['events'])} 件",
        "",
    ]

    by_stage: dict[int, list[dict]] = {}
    for e in d["events"]:
        by_stage.setdefault(e["stage_no"], []).append(e)

    for s in d["stages"]:
        no = s["stage_no"]
        if stage_no is not None and no != stage_no:
            continue
        out += [f"## {no}. {s['title']}", ""]
        if s.get("summary"):
            out += [s["summary"], ""]
        evs = by_stage.get(no, [])
        if not evs:
            out += ["（この段階に紐付く記録はありません）", ""]
            continue
        for e in evs:
            out += [
                f"### {e['occurred_on']}　{e['summary']}",
                "",
                f"- 種別: {e['kind']}　/　確度: {e['certainty']}",
                f"- 出典: {e['doc_title']}（{e['source_path']}）",
                "",
                "> " + (e["span_quote"] or "").replace("\n", "\n> "),
                "",
            ]

    if d["gaps"] and stage_no is None:
        out += ["## 記録の空白", ""]
        out += [f"- {g['description']}" for g in d["gaps"]] + [""]

    return "\n".join(out).encode("utf-8")


# ---------------------------------------------------------------- pptx

def build_pptx(ledger: Ledger, thread_id: str, stage_no: int | None = None) -> bytes:
    """段階ごとに1スライド。ストーリーテリングの流れをそのまま資料化する。"""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    d = _collect(ledger, thread_id, stage_no)
    th = d["thread"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    by_stage: dict[int, list[dict]] = {}
    for e in d["events"]:
        by_stage.setdefault(e["stage_no"], []).append(e)

    # 表紙
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2)).text_frame
    tb.text = th["name"]
    tb.paragraphs[0].runs[0].font.size = Pt(40)
    tb.paragraphs[0].runs[0].font.bold = True
    p = tb.add_paragraph()
    p.text = f"{th['first_seen']} 〜 {th['last_seen']}　／　記録 {len(d['events'])} 件"
    p.runs[0].font.size = Pt(16)
    p.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p2 = tb.add_paragraph()
    p2.text = f"出力日 {d['generated_at']}　業務確定台帳より自動生成"
    p2.runs[0].font.size = Pt(11)
    p2.runs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # 段階スライド
    for st in d["stages"]:
        no = st["stage_no"]
        if stage_no is not None and no != stage_no:
            continue
        s = prs.slides.add_slide(prs.slide_layouts[6])
        rgb = RGBColor(*_STAGE_RGB.get(no, (0x33, 0x33, 0x33)))

        bar = s.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.9)).text_frame
        bar.text = f"{no}. {st['title']}"
        bar.paragraphs[0].runs[0].font.size = Pt(30)
        bar.paragraphs[0].runs[0].font.bold = True
        bar.paragraphs[0].runs[0].font.color.rgb = rgb

        if st.get("summary"):
            sm = s.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(0.8)).text_frame
            sm.word_wrap = True
            sm.text = st["summary"]
            sm.paragraphs[0].runs[0].font.size = Pt(14)
            sm.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        body = s.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.7)).text_frame
        body.word_wrap = True
        evs = by_stage.get(no, [])
        if not evs:
            body.text = "この段階に紐付く記録はありません"
            body.paragraphs[0].runs[0].font.size = Pt(13)
            body.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        else:
            for i, e in enumerate(evs[:8]):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = f"{e['occurred_on']}　{e['summary']}"
                para.runs[0].font.size = Pt(14)
                para.space_after = Pt(6)
            if len(evs) > 8:
                more = body.add_paragraph()
                more.text = f"…ほか {len(evs) - 8} 件"
                more.runs[0].font.size = Pt(11)
                more.runs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------- xlsx

def build_xlsx(ledger: Ledger, thread_id: str, stage_no: int | None = None) -> bytes:
    """時系列の表と、段階別件数の棒グラフ（Enhancement.md §2-2 のグラフ/表）。"""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.styles import Alignment, Font, PatternFill

    d = _collect(ledger, thread_id, stage_no)
    th = d["thread"]
    wb = Workbook()

    ws = wb.active
    ws.title = "時系列"
    ws["A1"] = th["name"]
    ws["A1"].font = Font(size=14, bold=True)
    ws["A2"] = f"{th['first_seen']} 〜 {th['last_seen']}　／　出力日 {d['generated_at']}"
    ws["A2"].font = Font(size=9, color="666666")

    headers = ["段階", "段階名", "発生日", "種別", "内容", "確度", "出典", "原文"]
    ws.append([])
    ws.append(headers)
    hrow = ws.max_row
    for c in range(1, len(headers) + 1):
        cell = ws.cell(row=hrow, column=c)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="374A67")
        cell.alignment = Alignment(horizontal="center")

    for e in d["events"]:
        ws.append([
            e["stage_no"],
            STAGE_TITLES.get(e["stage_no"], ""),
            e["occurred_on"],
            e["kind"],
            e["summary"],
            e["certainty"],
            e["doc_title"],
            (e["span_quote"] or "").replace("\n", " "),
        ])

    for col, w in zip("ABCDEFGH", (6, 16, 12, 16, 46, 8, 28, 60)):
        ws.column_dimensions[col].width = w
    ws.freeze_panes = ws.cell(row=hrow + 1, column=1)

    # 段階別集計 + グラフ
    ws2 = wb.create_sheet("段階別集計")
    ws2.append(["段階", "段階名", "件数"])
    for c in range(1, 4):
        cell = ws2.cell(row=1, column=c)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="374A67")

    counts: dict[int, int] = {}
    for e in d["events"]:
        counts[e["stage_no"]] = counts.get(e["stage_no"], 0) + 1
    for no in sorted(STAGE_TITLES):
        if stage_no is not None and no != stage_no:
            continue
        ws2.append([no, STAGE_TITLES[no], counts.get(no, 0)])

    last = ws2.max_row
    if last > 1:
        chart = BarChart()
        chart.title = "段階別の記録件数"
        chart.y_axis.title = "件数"
        chart.x_axis.title = "段階"
        chart.add_data(Reference(ws2, min_col=3, min_row=1, max_row=last), titles_from_data=True)
        chart.set_categories(Reference(ws2, min_col=2, min_row=2, max_row=last))
        chart.height, chart.width = 8, 16
        ws2.add_chart(chart, "E2")
    ws2.column_dimensions["A"].width = 6
    ws2.column_dimensions["B"].width = 20
    ws2.column_dimensions["C"].width = 8

    if d["gaps"]:
        ws3 = wb.create_sheet("記録の空白")
        ws3.append(["種別", "起点", "内容"])
        for c in range(1, 4):
            cell = ws3.cell(row=1, column=c)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="B02A2A")
        for g in d["gaps"]:
            ws3.append([g["kind"], g["period_start"], g["description"]])
        for col, w in zip("ABC", (22, 14, 70)):
            ws3.column_dimensions[col].width = w

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


BUILDERS = {"md": build_md, "pptx": build_pptx, "xlsx": build_xlsx}
MEDIA_TYPES = {
    "md": "text/markdown; charset=utf-8",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


def build(ledger: Ledger, thread_id: str, fmt: str, stage_no: int | None = None) -> tuple[bytes, str, str]:
    """(内容, ファイル名, media_type) を返す。"""
    if fmt not in BUILDERS:
        raise ValueError(f"unsupported format: {fmt}")
    data = BUILDERS[fmt](ledger, thread_id, stage_no)
    name = filename(_thread(ledger, thread_id)["name"], fmt, stage_no)
    return data, name, MEDIA_TYPES[fmt]


# ==========================================================================
# Phase 8-3: 会議用資料（pdf / word）と資料用画像入り pptx を追加登録
# 上の定義は一切変更していない。以下はすべて追加。
# ==========================================================================

from projecttree import docs as _docs  # noqa: E402

BUILDERS["pdf"] = _docs.build_pdf
BUILDERS["docx"] = _docs.build_docx
BUILDERS["pptx_img"] = _docs.build_pptx_slides

MEDIA_TYPES["pdf"] = "application/pdf"
MEDIA_TYPES["docx"] = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MEDIA_TYPES["pptx_img"] = MEDIA_TYPES["pptx"]

# 形式キーと実際の拡張子が違うものだけを持つ。無い場合はキーがそのまま拡張子。
FILE_EXT = {"pptx_img": "pptx"}

FORMAT_LABELS = {
    "md": "md（台帳の時系列と原文のまま）",
    "pdf": "pdf（会議用資料・画像入り）",
    "docx": "word（会議用資料・画像入り）",
    "pptx": "pptx（プレゼン資料）",
    "pptx_img": "pptx（資料用画像スライド）",
    "xlsx": "excel（表・グラフ）",
}


def build_ext(ledger: Ledger, thread_id: str, fmt: str,
              stage_no: int | None = None) -> tuple[bytes, str, str]:
    """build() と同じだが、拡張子が形式キーと違う形式にも対応する。"""
    if fmt not in BUILDERS:
        raise ValueError(f"unsupported format: {fmt}")
    data = BUILDERS[fmt](ledger, thread_id, stage_no)
    name = filename(_thread(ledger, thread_id)["name"], FILE_EXT.get(fmt, fmt), stage_no)
    return data, name, MEDIA_TYPES[fmt]
